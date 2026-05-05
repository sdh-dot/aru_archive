"""Aru Archive v0.6.3 User Guide PPT generator."""
from __future__ import annotations
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
IMGS = Path(r"D:\AruArchive-v0.6.3-win-x64\매뉴얼용_스크린샷")
OUT  = Path(r"D:\AruArchive-v0.6.3-win-x64\AruArchive_User_Guide_v0.6.3.pptx")

SW = Inches(13.333)   # slide width
SH = Inches(7.5)      # slide height

# Colour palette (dark theme)
C_BG       = RGBColor(0x12, 0x12, 0x1E)   # near-black navy
C_TITLE    = RGBColor(0x64, 0xB4, 0xFF)   # sky-blue
C_BODY     = RGBColor(0xDC, 0xDC, 0xE6)   # light grey
C_ACCENT   = RGBColor(0x50, 0xC8, 0xB4)   # teal
C_WARN     = RGBColor(0xFF, 0xC8, 0x3C)   # amber
C_DIM      = RGBColor(0x88, 0x88, 0xA0)   # muted
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL    = RGBColor(0x1E, 0x1E, 0x32)   # slightly lighter than bg

HEADER_H   = Inches(0.95)   # title bar height
MARGIN     = Inches(0.35)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank_layout = prs.slide_layouts[6]   # blank


def slide() -> object:
    s = prs.slides.add_slide(blank_layout)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return s


def box(s, x, y, w, h, text="", size=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, bg=None, wrap=True):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_BODY
    if bg:
        tf.auto_size = None
        fill = tf._txBody.getparent().fill
        fill.solid()
        fill.fore_color.rgb = bg
    return tf


def title_bar(s, text, sub=None):
    """Dark accent bar across the top."""
    bar = s.shapes.add_shape(1, 0, 0, SW, HEADER_H)   # 1=rectangle
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PANEL
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = C_TITLE
    if sub:
        box(s, Inches(0.35), Inches(0.62), Inches(12), Inches(0.35),
            sub, size=13, color=C_DIM)


def img(s, fname, x, y, w, h=None):
    """Insert image; h=None → auto from aspect ratio."""
    p = IMGS / fname
    if not p.exists():
        return
    if h is None:
        from PIL import Image as _PI
        im = _PI.open(p)
        h = int(w * im.height / im.width)
        im.close()
    s.shapes.add_picture(str(p), x, y, w, h)
    return h


def label(s, x, y, text, color=C_ACCENT, size=12):
    box(s, x, y, Inches(4), Inches(0.35), text, size=size,
        bold=True, color=color)


def note_box(s, x, y, w, text, icon="💡", color=C_WARN):
    """Tip / warning callout box."""
    shape = s.shapes.add_shape(1, x, y, w, Inches(0.45))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x2A, 0x26, 0x10)
    shape.line.color.rgb = color
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{icon}  {text}"
    r.font.size = Pt(12); r.font.color.rgb = color


def bullet_tf(s, x, y, w, h, items, size=15, indent="  • "):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        r = p.add_run()
        r.text = f"{indent}{item}"
        r.font.size = Pt(size); r.font.color.rgb = C_BODY


def step_flow(s, steps, x, y, col_w, row_h):
    """Horizontal step boxes."""
    for i, (num, text) in enumerate(steps):
        bx = x + i * (col_w + Inches(0.12))
        sh = s.shapes.add_shape(1, bx, y, col_w, row_h)
        sh.fill.solid(); sh.fill.fore_color.rgb = C_PANEL
        sh.line.color.rgb = C_ACCENT
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"{'⓪①②③④⑤⑥⑦⑧⑨'[i]}  {num}"
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = C_TITLE
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = text
        r2.font.size = Pt(11); r2.font.color.rgb = C_BODY


# ===========================================================================
# Slide 1 — Cover
# ===========================================================================
s = slide()
# Logo/accent bar at top
top = s.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
top.fill.solid(); top.fill.fore_color.rgb = C_ACCENT
top.line.fill.background()

box(s, Inches(1.2), Inches(1.6), Inches(11), Inches(1.4),
    "Aru Archive", size=60, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
box(s, Inches(1.2), Inches(2.9), Inches(11), Inches(0.7),
    "v0.6.3  사용자 가이드", size=32, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(1.2), Inches(3.55), Inches(11), Inches(0.5),
    "Windows Release 기준  |  초보 사용자를 위한 설치부터 분류 실행까지",
    size=16, color=C_DIM, align=PP_ALIGN.CENTER)

# bottom bar
bot = s.shapes.add_shape(1, 0, SH - Inches(0.6), SW, Inches(0.6))
bot.fill.solid(); bot.fill.fore_color.rgb = C_PANEL
bot.line.fill.background()
box(s, Inches(0.4), SH - Inches(0.55), Inches(12), Inches(0.5),
    "⚠  배포 전 내용을 최신 버전으로 확인하세요",
    size=12, color=C_DIM, align=PP_ALIGN.CENTER)


# ===========================================================================
# Slide 2 — 전체 사용 흐름
# ===========================================================================
s = slide()
title_bar(s, "전체 사용 흐름", "설치부터 분류 완료까지 — 순서대로 따라하세요")

steps_a = [
    ("설치", "ZIP 해제\n.exe 실행"),
    ("폴더 설정", "분류 대상/\n완료 폴더"),
    ("이미지 스캔", "신규 파일\n감지"),
    ("메타데이터\n확인", "DB 조회\n상태 확인"),
    ("중복 검사", "시각적\n중복 확인"),
]
steps_b = [
    ("메타데이터\n보강", "Pixiv 조회\n+ XMP 등록"),
    ("분류 기준\n선택", "시리즈/\n캐릭터 등"),
    ("분류\n미리보기", "경로 확인\n전 검토"),
    ("분류 실행", "복사본\n생성"),
    ("결과 확인", "정합성\n보고서"),
]
cw = Inches(2.3)
rh = Inches(1.35)
step_flow(s, steps_a, MARGIN, HEADER_H + Inches(0.35), cw, rh)
step_flow(s, steps_b, MARGIN, HEADER_H + Inches(0.35) + rh + Inches(0.18), cw, rh)

note_box(s, MARGIN, SH - Inches(0.8), Inches(12.6),
         "각 단계는 순서대로 진행하세요. 메타데이터 보강 전에 분류를 실행하면 태그 정보 누락으로 분류 오류가 발생할 수 있습니다.")


# ===========================================================================
# Slide 3 — 파일명 / 이미지 준비
# ===========================================================================
s = slide()
title_bar(s, "Step 0  |  파일명 / 이미지 준비",
          "Aru Archive가 인식하는 파일명 형식을 확인합니다")

# Two large screenshots side by side
iw = Inches(6.3)
ih = img(s, "00_file_name_01.png", MARGIN, HEADER_H + Inches(0.25), iw)
img(s, "00_file_name_02.png", MARGIN + iw + Inches(0.2), HEADER_H + Inches(0.25), iw)

# small 3rd screenshot + explanation
img(s, "00_file_name_03.png", MARGIN, HEADER_H + Inches(0.25) + Inches(2.2), Inches(3.0))

bullets = [
    "Pixiv 작품 ID가 파일명에 포함되어야 메타데이터 자동 인식 가능",
    "예: 12345678_p0.jpg  /  illust_12345678_20240101.png",
    "일부 다운로더는 자동으로 올바른 형식으로 저장합니다",
    "지원 확장자: jpg, jpeg, png, gif, webp",
]
bullet_tf(s, MARGIN + Inches(3.2), HEADER_H + Inches(0.25) + Inches(2.0),
          Inches(9.5), Inches(1.8), bullets, size=14)

note_box(s, MARGIN, SH - Inches(0.8), Inches(12.6),
         "파일명에 Pixiv ID가 없으면 메타데이터 가져오기가 불가능합니다. 이름을 바꾸거나 ID가 포함된 다운로더를 사용하세요.", icon="⚠", color=C_WARN)


# ===========================================================================
# Slide 4 — 첫 실행 폴더 설정
# ===========================================================================
s = slide()
title_bar(s, "Step 1  |  첫 실행 — 폴더 설정",
          "처음 실행 시 세 가지 폴더를 지정합니다")

iw = Inches(5.8)
img(s, "01_first_run_folder_setup.png", MARGIN, HEADER_H + Inches(0.2), iw)
img(s, "02_first_run.png", MARGIN + iw + Inches(0.2), HEADER_H + Inches(0.2), iw)

folders = [
    ("분류 대상 폴더", "Pixiv에서 다운로드한 이미지가 있는 폴더"),
    ("분류 완료 폴더", "Classified 복사본이 저장될 폴더"),
    ("관리 폴더",      "DB · 설정이 저장되는 폴더 (기본: C:\\Users\\<user>\\AruArchive)"),
]
y = HEADER_H + Inches(0.25) + Inches(3.2)
for name, desc in folders:
    box(s, MARGIN, y, Inches(2.4), Inches(0.38), name, size=13, bold=True, color=C_ACCENT)
    box(s, MARGIN + Inches(2.5), y, Inches(10.3), Inches(0.38), desc, size=13, color=C_BODY)
    y += Inches(0.42)

note_box(s, MARGIN, SH - Inches(0.8), Inches(12.6),
         "관리 폴더는 한 번 설정 후 변경하지 않는 것을 권장합니다. 변경 시 DB를 수동으로 이전해야 합니다.", icon="⚠", color=C_WARN)


# ===========================================================================
# Slide 5 — 첫 단계 / 작업 폴더 확인
# ===========================================================================
s = slide()
title_bar(s, "Step 1  |  첫 단계 화면", "폴더 설정 완료 후 표시되는 메인 화면입니다")

iw = Inches(9.5)
ih = img(s, "03_first_step.png", (SW - iw) / 2, HEADER_H + Inches(0.2), iw)

bullets = [
    "좌측 사이드바: 현재 설정된 스캔/분류 폴더 경로 확인",
    "상단 탭: 각 단계(Step 1~9)로 이동",
    "'다음' 버튼으로 순서대로 진행",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.25) + Inches(3.6),
          Inches(12.6), Inches(1.2), bullets, size=14)


# ===========================================================================
# Slide 6 — 디렉토리 설정
# ===========================================================================
s = slide()
title_bar(s, "Step 2  |  디렉토리 설정", "스캔 대상 폴더와 분류 완료 폴더를 재확인합니다")

iw = Inches(9.5)
img(s, "04_디렉토리설정.png", (SW - iw) / 2, HEADER_H + Inches(0.2), iw)

bullets = [
    "하위 폴더 포함 여부 옵션을 확인하세요",
    "분류 완료 폴더는 스캔 대상 폴더와 다른 경로여야 합니다",
    "설정 변경 후 '다음'을 클릭하면 다음 단계로 넘어갑니다",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.25) + Inches(3.7),
          Inches(12.6), Inches(1.2), bullets, size=14)


# ===========================================================================
# Slide 7 — 이미지 스캔
# ===========================================================================
s = slide()
title_bar(s, "Step 3  |  이미지 스캔", "설정된 폴더에서 신규 이미지를 감지합니다")

iw = Inches(9.5)
img(s, "05_이미지_스캔.png", (SW - iw) / 2, HEADER_H + Inches(0.2), iw)

bullets = [
    "'스캔 시작' 클릭 → 대상 폴더를 순회해 신규 파일을 DB에 등록",
    "이미 등록된 파일은 건너뜁니다 (중복 스캔 안전)",
    "스캔 결과: 신규 N건 / 기존 M건 표시",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.25) + Inches(3.7),
          Inches(12.6), Inches(1.0), bullets, size=14)
note_box(s, MARGIN, SH - Inches(0.8), Inches(12.6),
         "스캔만으로는 파일이 이동/복사되지 않습니다. DB에 파일 경로를 등록하는 단계입니다.")


# ===========================================================================
# Slide 8 — 메타데이터 확인
# ===========================================================================
s = slide()
title_bar(s, "Step 4  |  메타데이터 확인", "스캔된 파일의 태그·시리즈 정보를 미리 확인합니다")

iw = Inches(9.5)
img(s, "06_메타데이터_확인.png", (SW - iw) / 2, HEADER_H + Inches(0.2), iw)

bullets = [
    "시리즈 · 캐릭터 · 작가 태그가 올바르게 인식됐는지 확인",
    "메타데이터 없음(노란색)은 다음 단계 '메타데이터 보강'으로 해결",
    "태그가 이미 있는 파일은 보강 없이 분류 단계로 직행 가능",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.25) + Inches(3.7),
          Inches(12.6), Inches(1.1), bullets, size=14)


# ===========================================================================
# Slide 9 — 시각적 중복 확인  (2×2 grid)
# ===========================================================================
s = slide()
title_bar(s, "Step 5  |  시각적 중복 확인", "비슷한 이미지를 나란히 보여주며 삭제 대상 선택")

iw = Inches(6.0)
ih_a = img(s, "07_시각적_중복확인_01.png", MARGIN, HEADER_H + Inches(0.15), iw)
img(s, "07_시각적_중복확인_02.png",
    MARGIN + iw + Inches(0.15), HEADER_H + Inches(0.15), iw)

row2_y = HEADER_H + Inches(0.15) + Inches(2.65)
img(s, "07_시각적_중복확인_03.png", MARGIN, row2_y, iw)
img(s, "07_시각적_중복확인_04.png", MARGIN + iw + Inches(0.15), row2_y, iw)

note_box(s, MARGIN, SH - Inches(0.78), Inches(12.6),
         "삭제는 직접 수행하지 않습니다. 중복으로 표시된 항목은 참고용이며 최종 판단은 사용자가 합니다.")


# ===========================================================================
# Slide 10 — 메타데이터 보강  (스크롤형 3장)
# ===========================================================================
s = slide()
title_bar(s, "Step 6  |  메타데이터 보강",
          "Pixiv 조회 → DB 저장 → XMP/탐색기 메타데이터 등록까지 한 번에 수행")

iw = Inches(4.15)
imgs_10 = ["08_메타데이터_보강_01.png", "08_메타데이터_보강_02.png", "08_메타데이터_보강_03.png"]
for i, fn in enumerate(imgs_10):
    img(s, fn, MARGIN + i * (iw + Inches(0.08)), HEADER_H + Inches(0.15), iw)

bullets = [
    "Pixiv 공개 API로 작품 정보(태그·시리즈·캐릭터)를 가져옵니다",
    "가져온 정보는 DB에 저장되고 파일의 XMP/Windows 태그에도 기록됩니다",
    "실패한 항목: 네트워크 오류 또는 비공개 작품 → 나중에 재시도 가능",
    "완료 후 '메타데이터 확인' 탭에서 결과를 재확인하세요",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.15) + Inches(3.35),
          Inches(12.6), Inches(1.6), bullets, size=13)

note_box(s, MARGIN, SH - Inches(0.78), Inches(12.6),
         "메타데이터 보강은 분류 전 반드시 완료해야 합니다. 누락 시 시리즈·캐릭터 폴더 분류가 되지 않습니다.", icon="⚠", color=C_WARN)


# ===========================================================================
# Slide 11 — 분류 기준 선택
# ===========================================================================
s = slide()
title_bar(s, "Step 7  |  분류 기준 선택", "출력 폴더 구조를 결정하는 핵심 단계")

iw = Inches(9.0)
img(s, "09_분류기준선택.png", (SW - iw) / 2, HEADER_H + Inches(0.18), iw)

options = [
    ("시리즈 + 캐릭터 폴더 (권장)", "BySeries/시리즈명/캐릭터명/ 구조로 정리"),
    ("시리즈 폴더만",               "BySeries/시리즈명/ — 캐릭터 폴더 없이"),
    ("작가명 기준",                 "ByAuthor/작가명/ — 시리즈 없는 작품에 유용"),
    ("개별 태그별",                 "ByTag/ — 추후 지원 예정"),
]
y = HEADER_H + Inches(0.18) + Inches(3.55)
for name, desc in options:
    box(s, MARGIN, y, Inches(3.2), Inches(0.36), name, size=13, bold=True, color=C_ACCENT)
    box(s, MARGIN + Inches(3.3), y, Inches(9.6), Inches(0.36), desc, size=13, color=C_BODY)
    y += Inches(0.39)


# ===========================================================================
# Slide 12 — 분류 미리보기  (2+2 grid)
# ===========================================================================
s = slide()
title_bar(s, "Step 7  |  분류 미리보기", "실행 전 복사될 경로를 미리 확인합니다")

iw = Inches(6.0)
ih_r = Inches(2.55)
img(s, "10_분류미리보기_01.png", MARGIN, HEADER_H + Inches(0.15), iw, ih_r)
img(s, "10_분류미리보기_02.png", MARGIN + iw + Inches(0.15), HEADER_H + Inches(0.15), iw, ih_r)
img(s, "10_분류미리보기_03.png", MARGIN, HEADER_H + Inches(0.15) + ih_r + Inches(0.1), iw, ih_r)
img(s, "10_분류미리보기_04.png", MARGIN + iw + Inches(0.15), HEADER_H + Inches(0.15) + ih_r + Inches(0.1), iw, ih_r)

note_box(s, MARGIN, SH - Inches(0.78), Inches(12.6),
         "미리보기에서 경로가 잘못됐다면 분류 기준 선택으로 돌아가 옵션을 변경하세요. 이 단계에서는 파일이 복사되지 않습니다.")


# ===========================================================================
# Slide 13 — 분류 실행
# ===========================================================================
s = slide()
title_bar(s, "Step 8  |  분류 실행", "미리보기 확인 후 실제 복사를 실행합니다")

iw = Inches(6.15)
img(s, "11_분류실행_01.png", MARGIN, HEADER_H + Inches(0.2), iw)
img(s, "11_분류실행_02.png", MARGIN + iw + Inches(0.15), HEADER_H + Inches(0.2), iw)

points = [
    "원본 파일은 이동·삭제하지 않습니다 — 복사본만 생성",
    "이미 분류된 항목은 '이미 분류됨'으로 자동 스킵 (재실행 안전)",
    "진행률 표시: 복사 N건 / 스킵 M건 / 실패 0건",
    "분류 완료 폴더에 BySeries/ (또는 선택한 구조) 폴더가 생성됩니다",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.2) + Inches(3.4),
          Inches(12.6), Inches(1.5), points, size=14)

note_box(s, MARGIN, SH - Inches(0.78), Inches(12.6),
         "실행 중 앱을 종료하지 마세요. 중단된 경우 재실행하면 이미 복사된 항목은 스킵하고 이어서 진행합니다.")


# ===========================================================================
# Slide 14 — 분류 결과 확인  (1 + 2×2)
# ===========================================================================
s = slide()
title_bar(s, "Step 9  |  분류 결과 확인", "정합성 보고서로 복사 누락·오류를 검증합니다")

# First image — full width top
iw_top = Inches(12.6)
ih_top = Inches(2.1)
img(s, "12_분류실행결과_01.png", MARGIN, HEADER_H + Inches(0.15), iw_top, ih_top)

iw_sm = Inches(6.15); ih_sm = Inches(2.0)
row2 = HEADER_H + Inches(0.15) + ih_top + Inches(0.1)
img(s, "12_분류실행결과_02.png", MARGIN,                      row2, iw_sm, ih_sm)
img(s, "12_분류실행결과_03.png", MARGIN + iw_sm + Inches(0.15), row2, iw_sm, ih_sm)

note_box(s, MARGIN, SH - Inches(0.78), Inches(12.6),
         "결과 화면에서 '실패 0건'을 반드시 확인하세요. 실패 항목은 원인 메시지를 확인 후 재실행하세요.")


# ===========================================================================
# Slide 14b — 분류 결과 (추가 스크린샷)
# ===========================================================================
s = slide()
title_bar(s, "Step 9  |  분류 결과 — 탐색기 / 갤러리 확인",
          "탐색기와 갤러리 뷰에서 분류된 파일을 확인합니다")

iw = Inches(6.15); ih = Inches(2.9)
img(s, "12_분류실행결과_04.png", MARGIN,                      HEADER_H + Inches(0.2), iw, ih)
img(s, "12_분류실행결과_05.png", MARGIN + iw + Inches(0.15), HEADER_H + Inches(0.2), iw, ih)

bullets = [
    "탐색기: BySeries/시리즈명/캐릭터명/ 경로에 파일 존재 확인",
    "갤러리 뷰: 앱 내에서 분류된 이미지를 카드 형태로 탐색 가능",
    "legacy_extra: 이전 분류 결과 중 현재 규칙과 맞지 않는 파일",
    "missing_expected: 예상된 복사본이 없는 항목 → 재실행 필요",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.2) + ih + Inches(0.15),
          Inches(12.6), Inches(1.5), bullets, size=14)


# ===========================================================================
# Slide 15 — 브라우저 확장 설치
# ===========================================================================
s = slide()
title_bar(s, "브라우저 확장 설치",
          "Chrome / Edge에서 Pixiv 작품을 바로 Aru Archive에 추가합니다")

iw_a = Inches(5.5)
iw_b = Inches(7.1)
img(s, "13_브라우저_확장_설치_01.png", MARGIN, HEADER_H + Inches(0.2), iw_a)
img(s, "14_브라우저_확장_설치_02.png", MARGIN + iw_a + Inches(0.15), HEADER_H + Inches(0.2), iw_b)

steps_ext = [
    "Chrome / Edge 주소창에  chrome://extensions  입력",
    "우측 상단 '개발자 모드' 토글 활성화",
    "'압축 해제된 확장 프로그램 로드' 클릭",
    "ZIP이 아닌 manifest.json이 있는 폴더 선택",
    "확장 아이콘 클릭 → Pixiv 페이지에서 다운로드 버튼 표시",
]
bullet_tf(s, MARGIN, HEADER_H + Inches(0.2) + Inches(3.55),
          Inches(12.6), Inches(2.0), steps_ext, size=14,
          indent="  ")

note_box(s, MARGIN, SH - Inches(0.78), Inches(12.6),
         "현재 Chrome/Edge 마켓에 미등록 상태입니다. 개발자 모드 로드 방식을 사용하세요.", icon="⚠", color=C_WARN)


# ===========================================================================
# Slide 16 — 자주 묻는 문제
# ===========================================================================
s = slide()
title_bar(s, "자주 묻는 문제 (FAQ)", "")

faqs = [
    ("SmartScreen 경고가 뜹니다",
     "'추가 정보' → '실행' 클릭. 마켓 미등록 앱이므로 정상입니다."),
    ("ZIP을 해제하지 않고 실행했습니다",
     "ZIP 안에서 실행하면 폴더 설정이 저장되지 않습니다. 반드시 압축 해제 후 실행하세요."),
    ("메타데이터 가져오기가 실패합니다",
     "비공개 작품이거나 Pixiv 점검 중일 수 있습니다. 잠시 후 재시도하세요."),
    ("분류 실행 후 '이미 분류됨'만 표시됩니다",
     "이전에 같은 파일을 이미 분류한 경우입니다. 정상 동작이며 중복 복사를 방지합니다."),
    ("legacy_extra 항목이 있습니다",
     "이전 분류 기준으로 복사된 파일입니다. 수동으로 삭제하거나 무시해도 됩니다."),
    ("missing_expected 항목이 있습니다",
     "예상 경로에 파일이 없습니다. 분류 실행을 다시 진행하면 복사됩니다."),
]

col_w = Inches(6.15)
rows_l = faqs[:3]
rows_r = faqs[3:]
y_base = HEADER_H + Inches(0.2)
for row_list, x_off in [(rows_l, MARGIN), (rows_r, MARGIN + col_w + Inches(0.15))]:
    y = y_base
    for q, a in row_list:
        qb = s.shapes.add_shape(1, x_off, y, col_w, Inches(0.32))
        qb.fill.solid(); qb.fill.fore_color.rgb = C_PANEL
        qb.line.fill.background()
        tf = qb.text_frame
        r = tf.paragraphs[0].add_run()
        r.text = f"Q  {q}"
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = C_ACCENT

        box(s, x_off, y + Inches(0.33), col_w, Inches(0.52),
            f"A  {a}", size=12, color=C_BODY)
        y += Inches(0.9)


# ===========================================================================
# Slide 17 — 마무리 체크리스트
# ===========================================================================
s = slide()
title_bar(s, "마무리 체크리스트", "분류 완료 전 이 항목들을 모두 확인하세요")

checks = [
    ("폴더 설정 완료",
     "분류 대상 / 분류 완료 / 관리 폴더가 올바르게 지정됐는지 확인"),
    ("이미지 스캔 완료",
     "신규 파일이 DB에 등록됐는지 — 스캔 결과에서 신규 N건 확인"),
    ("메타데이터 보강 완료",
     "메타데이터 없음 항목이 0건이 될 때까지 보강"),
    ("분류 미리보기 확인",
     "예상 경로가 올바른지 실행 전 반드시 검토"),
    ("분류 실행 완료",
     "실패 0건 확인 — 복사 N건 / 스킵 M건 정상 표시"),
    ("정합성 보고서 확인",
     "missing_expected 0건 / legacy_extra 처리 방침 결정"),
]

x_left  = MARGIN
x_right = MARGIN + Inches(6.4)
y = HEADER_H + Inches(0.3)
for i, (title_c, desc) in enumerate(checks):
    x = x_left if i < 3 else x_right
    row_y = y + (i % 3) * Inches(1.78)

    cb = s.shapes.add_shape(1, x, row_y, Inches(5.9), Inches(1.55))
    cb.fill.solid(); cb.fill.fore_color.rgb = C_PANEL
    cb.line.color.rgb = C_ACCENT

    tf = cb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"☐  {title_c}"
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = C_TITLE
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = f"     {desc}"
    r2.font.size = Pt(12); r2.font.color.rgb = C_BODY

# Final note
note_box(s, MARGIN, SH - Inches(0.72), Inches(12.6),
         "모든 항목을 확인했다면 브라우저 확장을 설치해 Pixiv에서 바로 다운로드·분류 연동을 시작하세요 🎉",
         icon="✅", color=C_ACCENT)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved → {OUT}")

# Screenshot index
used = []
import re
src = Path(__file__).read_text(encoding="utf-8")
for m in re.finditer(r'img\(s,\s*"([^"]+)"', src):
    fn = m.group(1)
    if fn not in used:
        used.append(fn)
print("\n사용된 스크린샷 목록:")
for u in used:
    print(f"  {u}")
not_used = [f.name for f in IMGS.glob("*.png") if f.name not in used]
if not_used:
    print("\n미사용 스크린샷:")
    for n in not_used:
        print(f"  {n}")
